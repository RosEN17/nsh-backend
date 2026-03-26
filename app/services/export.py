import io
from docx import Document
from typing import Tuple

def build_variance_analysis_markdown(pack, report_items, spec) -> str:
    lines = [
        f"# {spec.get('title', 'NordSheet Report')}",
        f"## Period",
        f"- Current: {pack.current_period}",
        f"- Previous: {pack.previous_period or 'N/A'}",
        "",
        "## Narrative",
        pack.narrative or "Ingen narrativ kommentar tillgänglig.",
        "",
        "## Saved report items",
    ]

    if report_items:
        for item in report_items:
            lines.append(f"- {item.get('text', '')}")
    else:
        lines.append("- Inga sparade rapportpunkter.")

    return "\n".join(lines)

def build_finance_report_deck(pack, report_items, spec, btype="Övrigt") -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = spec.get("title", "NordSheet Finance Report")
    slide.placeholders[1].text = f"Current period: {pack.current_period}\nPrevious period: {pack.previous_period or 'N/A'}"

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Narrative"
    slide2.placeholders[1].text = pack.narrative or "Ingen narrativ kommentar."

    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Saved Report Items"
    body = "\n".join([f"• {x.get('text', '')}" for x in report_items]) if report_items else "• Inga sparade rapportpunkter."
    slide3.placeholders[1].text = body

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

def export_artifact(fmt: str, pack, report_items, spec, purpose="report", btype="Övrigt") -> Tuple[bytes, str, str]:
    safe_period = str(pack.current_period).replace("/", "-")

    if fmt == "pptx":
        data = build_finance_report_deck(pack=pack, report_items=report_items, spec=spec, btype=btype)
        return (
            data,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            f"{purpose}_{safe_period}.pptx",
        )

    content = build_variance_analysis_markdown(pack, report_items, spec)

    if fmt == "docx":
        doc = Document()
        for line in content.splitlines():
            if line.startswith("# "):
                doc.add_heading(line.replace("# ", ""), level=1)
            elif line.startswith("## "):
                doc.add_heading(line.replace("## ", ""), level=2)
            elif line.startswith("- "):
                doc.add_paragraph(line.replace("- ", ""), style="List Bullet")
            elif line.strip():
                doc.add_paragraph(line)

        buf = io.BytesIO()
        doc.save(buf)
        return (
            buf.getvalue(),
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            f"{purpose}_{safe_period}.docx",
        )

    if fmt == "markdown":
        return content.encode("utf-8"), "text/markdown", f"{purpose}_{safe_period}.md"

    return content.encode("utf-8"), "text/plain", f"{purpose}_{safe_period}.txt"
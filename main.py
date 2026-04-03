from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
import pandas as pd
import numpy as np
import json, io, os, re
from typing import Any, Optional

app = FastAPI(title="NordSheet API")

app.add_middleware(
    CORSMiddleware,
    allow_origin_regex=r"https://.*nordsheet\.com|https://nsh-frontend.*\.vercel\.app|http://localhost:3000",
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# ═══════════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════════

def read_upload_bytes(b: bytes, filename: str) -> pd.DataFrame:
    ext = os.path.splitext(filename or "")[1].lower()
    na_vals = ["N/A","NA","n/a","na","NULL","null","None","none","NaN","nan","-","--",""]
    if ext == ".csv":
        return pd.read_csv(io.BytesIO(b), keep_default_na=True, na_values=na_vals)
    elif ext in {".xlsx", ".xls"}:
        return pd.read_excel(io.BytesIO(b), keep_default_na=True, na_values=na_vals)
    raise ValueError("Endast CSV, XLSX och XLS stöds.")


def suggest_columns(cols: list) -> dict:
    cols_lower = {c: c.lower() for c in cols}
    keywords = {
        "period":       ["period","month","månad","date","datum","year","år"],
        "account":      ["account","konto","kontonr","account_id","gl"],
        "account_name": ["account_name","kontonamn","name","description","benämning"],
        "actual":       ["actual","utfall","fact","verklig","bokfört","belopp"],
        "budget":       ["budget","plan","planned","bud"],
        "entity":       ["entity","bolag","company","legal","enhet"],
        "cost_center":  ["cost_center","costcenter","cc","kostnadsställe","resultatenhet"],
        "project":      ["project","projekt"],
    }
    return {
        field: [c for c, cl in cols_lower.items() if any(kw in cl for kw in kws)]
        for field, kws in keywords.items()
    }


def fmt_sek(n: float, decimals: int = 0) -> str:
    abs_n = abs(n)
    if abs_n >= 1_000_000:
        return f"{n/1_000_000:.{decimals}f} MSEK"
    if abs_n >= 1_000:
        return f"{n/1_000:.{decimals}f} tkr"
    return f"{n:.{decimals}f} SEK"


def safe_pct(a: float, b: float) -> Optional[float]:
    if b and b != 0:
        return (a - b) / abs(b)
    return None


# ═══════════════════════════════════════════════════════════════════
# PACK COMPUTATION
# ═══════════════════════════════════════════════════════════════════

def compute_pack(df: pd.DataFrame, mapping: dict) -> dict:
    rename = {v: k for k, v in mapping.items() if v and v in df.columns}
    df = df.rename(columns=rename)
    warnings = []

    # Numeric coercion
    for col in ["actual", "budget"]:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0)

    # Periods
    periods = []
    if "period" in df.columns:
        periods = sorted(df["period"].dropna().astype(str).unique().tolist())
    current_period  = periods[-1] if periods else "Unknown"
    previous_period = periods[-2] if len(periods) >= 2 else None

    total_actual = float(df["actual"].sum()) if "actual" in df.columns else 0
    total_budget = float(df["budget"].sum()) if "budget" in df.columns else 0

    top_budget, top_mom = [], []
    kpi_summary = {}
    by_account = pd.DataFrame()

    if "account" in df.columns and "actual" in df.columns:
        name_col = "account_name" if "account_name" in df.columns else None

        agg_cols = {"actual": "sum"}
        if "budget" in df.columns:
            agg_cols["budget"] = "sum"

        grp_cols = ["account"]
        if name_col:
            grp_cols.append(name_col)

        by_account = df.groupby(grp_cols, dropna=False).agg(agg_cols).reset_index()

        if "budget" in by_account.columns:
            by_account["variance"]     = by_account["actual"] - by_account["budget"]
            by_account["variance_pct"] = by_account.apply(
                lambda r: safe_pct(r["actual"], r["budget"]), axis=1
            )
        else:
            by_account["variance"]     = 0
            by_account["variance_pct"] = None
            warnings.append("Budgetkolumn saknas — variansanalys ej möjlig.")

        def row_to_dict(r):
            label = str(r.get("account_name") or r.get("account") or "—")
            return {
                "Konto":         str(r.get("account", "—")),
                "Label":         label,
                "Utfall":        round(float(r["actual"]), 0),
                "Budget":        round(float(r.get("budget", 0)), 0),
                "Vs budget diff": round(float(r.get("variance", 0)), 0),
                "Vs budget %":   round(float(r.get("variance_pct") or 0), 4),
            }

        top_budget = [row_to_dict(r) for _, r in by_account.sort_values("variance").head(10).iterrows()]
        top_mom    = [row_to_dict(r) for _, r in by_account.sort_values("variance", ascending=False).head(10).iterrows()]

        # MoM: compare current vs previous period if both exist
        mom_diff, mom_pct = None, None
        if previous_period and "period" in df.columns:
            cur_total  = float(df[df["period"].astype(str) == current_period]["actual"].sum())
            prev_total = float(df[df["period"].astype(str) == previous_period]["actual"].sum())
            mom_diff   = cur_total - prev_total
            mom_pct    = safe_pct(cur_total, prev_total)

        kpi_summary = {
            "Nu":             round(total_actual, 0),
            "Föregående":     round(float(df[df["period"].astype(str) == previous_period]["actual"].sum()), 0) if previous_period and "period" in df.columns else None,
            "MoM diff":       round(mom_diff, 0) if mom_diff is not None else None,
            "MoM %":          round(mom_pct, 4) if mom_pct is not None else None,
            "Budget":         round(total_budget, 0),
            "Vs budget diff": round(total_actual - total_budget, 0),
            "Vs budget %":    round(safe_pct(total_actual, total_budget) or 0, 4),
        }

    # Period-level time series (for forecast charts)
    period_series = []
    if "period" in df.columns and "actual" in df.columns:
        ps = df.groupby("period").agg(
            actual=("actual", "sum"),
            **({"budget": ("budget", "sum")} if "budget" in df.columns else {})
        ).reset_index().sort_values("period")
        period_series = ps.fillna(0).to_dict(orient="records")

    # Narrative
    var = total_actual - total_budget
    var_pct = safe_pct(total_actual, total_budget)
    narrative = (
        f"Period {current_period}: Utfall {fmt_sek(total_actual)} "
        f"mot budget {fmt_sek(total_budget)} "
        f"(avvikelse {fmt_sek(var, 0)}"
        f"{f', {var_pct*100:+.1f}%' if var_pct is not None else ''})."
    )
    if previous_period and kpi_summary.get("MoM diff") is not None:
        narrative += (
            f" Jämfört med {previous_period}: "
            f"{fmt_sek(kpi_summary['MoM diff'])}"
            f" ({kpi_summary.get('MoM %', 0)*100:+.1f}%)."
        )

    return {
        "current_period":  current_period,
        "previous_period": previous_period,
        "periods":         periods,
        "warnings":        warnings,
        "narrative":       narrative,
        "top_budget":      top_budget,
        "top_mom":         top_mom,
        "kpi_summary":     [kpi_summary],
        "total_actual":    round(total_actual, 0),
        "total_budget":    round(total_budget, 0),
        "period_series":   period_series,
        "account_rows":    by_account.fillna(0).to_dict(orient="records") if not by_account.empty else [],
    }


# ═══════════════════════════════════════════════════════════════════
# FORECAST ENGINE
# ═══════════════════════════════════════════════════════════════════

def build_forecast_data(pack: dict) -> dict:
    """
    Build forecast from pack data.
    Uses linear regression on period_series for trend,
    then extrapolates forward with confidence intervals.
    """
    period_series = pack.get("period_series", [])
    total_actual  = pack.get("total_actual", 0)
    total_budget  = pack.get("total_budget", 0)
    kpi           = pack.get("kpi_summary", [{}])[0]
    mom_pct       = float(kpi.get("MoM %") or 0)

    # Build actuals array
    actuals = [float(p.get("actual", 0)) for p in period_series]
    budgets = [float(p.get("budget", 0)) for p in period_series]
    labels  = [str(p.get("period", "")) for p in period_series]

    if len(actuals) < 2:
        # Not enough data — use simple MoM extrapolation
        actuals = [total_actual * (0.9 + i * 0.02) for i in range(6)]
        budgets = [total_budget * (0.9 + i * 0.02) for i in range(6)]
        labels  = [f"P{i+1}" for i in range(6)]

    # Linear trend on actuals
    x = np.arange(len(actuals))
    if len(x) >= 2:
        coeffs    = np.polyfit(x, actuals, 1)
        slope     = float(coeffs[0])
        intercept = float(coeffs[1])
    else:
        slope, intercept = 0, actuals[-1] if actuals else 0

    # Std dev for confidence bands
    residuals = [a - (slope * i + intercept) for i, a in enumerate(actuals)]
    std_dev   = float(np.std(residuals)) if len(residuals) > 1 else abs(slope) * 0.15

    # Forecast 12 months forward
    n_hist    = len(actuals)
    fc_months = []
    for i in range(12):
        idx      = n_hist + i
        base_fc  = slope * idx + intercept
        growth   = base_fc * (1 + mom_pct * 0.3)  # damped MoM
        fc_months.append({
            "label":       f"F+{i+1}",
            "forecast":    round(max(growth, 0), 0),
            "upper":       round(max(growth + 1.96 * std_dev, 0), 0),
            "lower":       round(max(growth - 1.96 * std_dev, 0), 0),
        })

    # Scenarios
    last_fc = fc_months[-1]["forecast"] if fc_months else total_actual
    monthly_base = fc_months[0]["forecast"] if fc_months else total_actual

    scenarios = {
        "optimistic":  round(monthly_base * 1.07, 0),
        "base":        round(monthly_base * 1.03, 0),
        "pessimistic": round(monthly_base * 0.97, 0),
    }

    # Income/cost split (use account rows if available)
    account_rows = pack.get("account_rows", [])
    income_total = 0.0
    cost_total   = 0.0

    for row in account_rows:
        val = float(row.get("actual", 0))
        # Heuristic: positive = income, negative = cost
        if val >= 0:
            income_total += val
        else:
            cost_total += abs(val)

    if income_total == 0 and cost_total == 0:
        income_total = total_actual * 0.6
        cost_total   = total_actual * 0.4

    income_forecast = income_total * (1 + mom_pct)
    cost_forecast   = cost_total   * (1 + mom_pct * 0.7)
    margin          = ((income_forecast - cost_forecast) / income_forecast * 100) if income_forecast else 0

    return {
        "history":          [{"label": l, "actual": a, "budget": b} for l, a, b in zip(labels, actuals, budgets)],
        "forecast":         fc_months,
        "scenarios":        scenarios,
        "income_forecast":  round(income_forecast, 0),
        "cost_forecast":    round(cost_forecast, 0),
        "margin_pct":       round(margin, 1),
        "growth_rate":      round(mom_pct * 100, 1),
        "trend_slope":      round(slope, 0),
    }



# ═══════════════════════════════════════════════════════════════════
# PPTX BUILDER  — skräddarsydd per rapporttyp
# ═══════════════════════════════════════════════════════════════════

def build_pptx(pack: dict, spec: dict, ai_context: str = "") -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # ── Palette ──────────────────────────────────────────────────────────────
    C_BG      = RGBColor(0x0D, 0x0D, 0x12)
    C_CARD    = RGBColor(0x0F, 0x0F, 0x16)
    C_CARD2   = RGBColor(0x13, 0x13, 0x1E)
    C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
    C_MUTED   = RGBColor(0xA0, 0xA0, 0xB8)
    C_FAINT   = RGBColor(0x44, 0x44, 0x5A)
    C_ACCENT  = RGBColor(0x6C, 0x63, 0xFF)
    C_ACCENT2 = RGBColor(0x9B, 0x94, 0xFF)
    C_GREEN   = RGBColor(0x22, 0xC5, 0x5E)
    C_GREEN_D = RGBColor(0x16, 0x7A, 0x3B)
    C_RED     = RGBColor(0xEF, 0x44, 0x44)
    C_RED_D   = RGBColor(0x8B, 0x1A, 0x1A)
    C_AMBER   = RGBColor(0xF5, 0x9E, 0x0B)

    # ── Unpack data ──────────────────────────────────────────────────────────
    current_period  = pack.get("current_period",  "—")
    previous_period = pack.get("previous_period", "—")
    periods         = pack.get("periods", [])
    total_actual    = float(pack.get("total_actual", 0))
    total_budget    = float(pack.get("total_budget", 0))
    variance        = total_actual - total_budget
    var_pct         = safe_pct(total_actual, total_budget) or 0
    kpi             = pack.get("kpi_summary", [{}])[0]
    mom_pct         = float(kpi.get("MoM %") or 0)
    mom_diff        = float(kpi.get("MoM diff") or 0)
    narrative       = pack.get("narrative", "")
    top_budget      = pack.get("top_budget", [])[:10]
    top_mom         = pack.get("top_mom",    [])[:10]
    period_series   = pack.get("period_series", [])
    report_type     = spec.get("report_type",  "monthly")
    tone            = spec.get("tone",         "professionell")
    context_text    = spec.get("context",      "") or ai_context
    detail_level    = int(spec.get("detail",   60))
    fc_data         = build_forecast_data(pack)

    REPORT_TITLES = {
        "monthly":       "Månadsrapport",
        "quarterly":     "Kvartalsrapport",
        "annual":        "Årsbokslut / Årsredovisning",
        "forecast":      "Prognos & Forecast",
        "kpi_dashboard": "KPI-rapport",
        # legacy compat
        "income_statement": "Resultaträkning",
        "balance_sheet":    "Balansräkning",
        "cash_flow":        "Kassaflöde",
        "budget_vs_actual": "Budget vs Utfall",
        "ai_summary":       "AI-sammanfattning",
    }
    report_title = spec.get("title") or REPORT_TITLES.get(report_type, "Finansiell rapport")

    # ── AI — generate tailored content ───────────────────────────────────────
    ai_insights    = []
    ai_bullets     = {}        # keyed sections → list of bullet strings
    ai_exec_summ   = ""
    openai_key     = os.getenv("OPENAI_API_KEY")

    top3_txt = ", ".join([
        f"{r.get('Label', r.get('Konto','?'))}: {fmt_sek(r.get('Vs budget diff', 0))}"
        for r in top_budget[:3]
    ])
    fc = fc_data

    if openai_key:
        try:
            from openai import OpenAI
            client = OpenAI(api_key=openai_key)

            section_instructions = {
                "monthly": (
                    "Skapa innehåll för en professionell månadsrapport. "
                    "Returnera JSON med nycklarna: "
                    "executive_summary (2 meningar), "
                    "variance_drivers (lista med 4 punkter om avvikelsedrivare), "
                    "highlights (lista med 3 positiva punkter), "
                    "risks (lista med 3 risker eller åtgärder), "
                    "outlook (1-2 meningar om nästa period)."
                ),
                "quarterly": (
                    "Skapa innehåll för en professionell kvartalsrapport till ledningsgrupp. "
                    "Returnera JSON med nycklarna: "
                    "executive_summary (3 meningar med kvartalsperspektiv), "
                    "quarter_analysis (lista med 4 punkter om kvartalets prestanda), "
                    "budget_commentary (lista med 3 punkter om budget vs utfall), "
                    "strategic_actions (lista med 3 strategiska åtgärder), "
                    "next_quarter_outlook (2 meningar om nästa kvartal)."
                ),
                "annual": (
                    "Skapa innehåll för en professionell årsredovisning / styrelserapport. "
                    "Returnera JSON med nycklarna: "
                    "executive_summary (3 meningar som VD-kommentar), "
                    "year_highlights (lista med 4 punkter om årets prestanda), "
                    "financial_position (lista med 3 punkter om finansiell ställning), "
                    "board_recommendations (lista med 3 styrelserekommendationer), "
                    "next_year_outlook (2 meningar om kommande år)."
                ),
                "forecast": (
                    "Skapa innehåll för en finansiell prognos- och forecastrapport. "
                    "Returnera JSON med nycklarna: "
                    "executive_summary (2 meningar om prognosen), "
                    "assumptions (lista med 3 antaganden bakom prognosen), "
                    "scenario_commentary (lista med 3 punkter: optimistisk, bas, pessimistisk), "
                    "risk_factors (lista med 3 riskfaktorer), "
                    "recommendations (lista med 2 rekommendationer)."
                ),
                "kpi_dashboard": (
                    "Skapa innehåll för en KPI-rapport / dashboard-presentation. "
                    "Returnera JSON med nycklarna: "
                    "executive_summary (2 meningar om KPI-status), "
                    "kpi_highlights (lista med 4 KPI-observationer), "
                    "green_flags (lista med 2 positiva KPI:er), "
                    "red_flags (lista med 2 KPI:er som kräver åtgärd), "
                    "action_items (lista med 3 konkreta åtgärder)."
                ),
            }
            instr = section_instructions.get(
                report_type,
                section_instructions["monthly"]
            )
            prompt = (
                f"Du är en senior finansanalytiker. {instr}\n\n"
                f"Ton: {tone}. Allt på svenska.\n"
                f"Data:\n"
                f"- Period: {current_period} (föregående: {previous_period})\n"
                f"- Utfall: {fmt_sek(total_actual)}\n"
                f"- Budget: {fmt_sek(total_budget)}\n"
                f"- Avvikelse: {fmt_sek(variance)} ({var_pct*100:+.1f}%)\n"
                f"- MoM: {mom_pct*100:+.1f}% ({fmt_sek(mom_diff)})\n"
                f"- Topp avvikelser: {top3_txt}\n"
                f"- Prognos nästa period (bas): {fmt_sek(fc['scenarios']['base'])}\n"
                + (f"- Extra kontext: {context_text}\n" if context_text else "")
                + "\nReturnera ENDAST giltig JSON."
            )
            resp = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "Du är en elite finansrapportskribent. Returnera alltid JSON."},
                    {"role": "user",   "content": prompt},
                ],
                max_tokens=900,
            )
            raw = re.sub(r"```json|```", "", resp.choices[0].message.content or "{}").strip()
            ai_bullets     = json.loads(raw)
            ai_exec_summ   = ai_bullets.get("executive_summary", narrative)
            # Build flat insights for the insights slide
            ai_insights = []
            for k, v in ai_bullets.items():
                if k == "executive_summary":
                    continue
                if isinstance(v, list):
                    ai_insights.append({"title": k.replace("_", " ").capitalize(), "items": v})
                elif isinstance(v, str):
                    ai_insights.append({"title": k.replace("_", " ").capitalize(), "items": [v]})
        except Exception:
            pass

    if not ai_exec_summ:
        ai_exec_summ = narrative
    if not ai_insights:
        ai_insights = [
            {"title": "Utfall vs budget",    "items": [f"Utfall {fmt_sek(total_actual)} mot budget {fmt_sek(total_budget)}.", f"Avvikelse {fmt_sek(variance)} ({var_pct*100:+.1f}%)."]},
            {"title": "MoM-trend",           "items": [f"Förändring mot föregående period: {mom_pct*100:+.1f}%.", f"Absolut förändring: {fmt_sek(mom_diff)}."]},
            {"title": "Avvikelsedrivare",    "items": [f"{r.get('Label','?')}: {fmt_sek(r.get('Vs budget diff',0))}" for r in top_budget[:3]]},
            {"title": "Prognos",             "items": [f"Bas: {fmt_sek(fc['scenarios']['base'])}", f"Optimistisk: {fmt_sek(fc['scenarios']['optimistic'])}", f"Pessimistisk: {fmt_sek(fc['scenarios']['pessimistic'])}"]},
        ]

    # ── Presentation setup ───────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── Low-level drawing helpers ─────────────────────────────────────────────
    def add_rect(slide, x, y, w, h, color: RGBColor):
        s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = color
        s.line.fill.background()
        return s

    def add_text(slide, text, x, y, w, h, size=14, bold=False,
                 color=C_WHITE, align=PP_ALIGN.LEFT, italic=False):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        p  = tf.paragraphs[0]; p.alignment = align
        r  = p.add_run(); r.text = str(text)
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.italic = italic
        return tb

    def add_bullets(slide, items, x, y, w, h, size=11, color=C_MUTED):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = f"• {item}"
            r.font.size = Pt(size); r.font.color.rgb = color

    def bg(slide):
        add_rect(slide, 0, 0, 13.33, 7.5, C_BG)

    def accent_bar(slide, color=C_ACCENT):
        add_rect(slide, 0, 0, 0.07, 7.5, color)

    def slide_header(slide, title, subtitle="", tag=""):
        accent_bar(slide)
        if tag:
            add_text(slide, tag, 0.25, 0.22, 5, 0.28, size=8,
                     color=C_FAINT, bold=True)
        add_text(slide, title,    0.25, 0.45, 9, 0.55, size=24, bold=True)
        if subtitle:
            add_text(slide, subtitle, 0.25, 1.0, 9, 0.35, size=11, color=C_MUTED)

    def kpi_box(slide, x, y, w, h, label, value, sub="", val_color=C_WHITE):
        add_rect(slide, x, y, w, h, C_CARD)
        add_text(slide, label, x+0.14, y+0.10, w-0.28, 0.26, size=9,
                 color=C_FAINT, bold=True)
        add_text(slide, value, x+0.14, y+0.36, w-0.28, 0.50, size=19,
                 bold=True, color=val_color)
        if sub:
            add_text(slide, sub, x+0.14, y+h-0.35, w-0.28, 0.28, size=9, color=C_MUTED)

    def divider(slide, y):
        add_rect(slide, 0.25, y, 12.8, 0.015, C_FAINT)

    def bar_chart(slide, data, x, y, w, h, show_budget=True):
        """Render a grouped bar chart (actual vs budget) using rectangles."""
        if not data: return
        vals = []
        for d in data:
            vals.append(abs(float(d.get("actual", d.get("Utfall", 0)) or 0)))
            if show_budget:
                vals.append(abs(float(d.get("budget", d.get("Budget", 0)) or 0)))
        max_v = max(vals) if vals else 1
        if max_v == 0: max_v = 1

        n       = len(data)
        bar_w   = w / n
        area_h  = h - 0.45
        gap     = 0.03

        for i, d in enumerate(data):
            bx  = x + i * bar_w
            act = abs(float(d.get("actual", d.get("Utfall", 0)) or 0))
            bud = abs(float(d.get("budget", d.get("Budget", 0)) or 0)) if show_budget else 0

            if show_budget and bud > 0:
                bh = (bud / max_v) * area_h * 0.88
                add_rect(slide, bx + bar_w * 0.08, y + area_h - bh,
                         bar_w * 0.38, bh, C_CARD2)

            ah = (act / max_v) * area_h * 0.88 if act > 0 else 0.02
            bx_act = bx + (bar_w * 0.52 if show_budget else bar_w * 0.18)
            bw_act = bar_w * (0.38 if show_budget else 0.64)
            add_rect(slide, bx_act, y + area_h - ah, bw_act, ah, C_ACCENT)

            lbl = str(d.get("Label", d.get("label", d.get("Konto", ""))))[:9]
            add_text(slide, lbl, bx, y + area_h + 0.04, bar_w, 0.24,
                     size=7, color=C_FAINT, align=PP_ALIGN.CENTER)

        # Legend
        if show_budget:
            add_rect(slide, x, y + h - 0.28, 0.18, 0.12, C_ACCENT)
            add_text(slide, "Utfall",  x + 0.22, y + h - 0.32, 1.2, 0.22, size=8, color=C_MUTED)
            add_rect(slide, x + 1.2,  y + h - 0.28, 0.18, 0.12, C_CARD2)
            add_text(slide, "Budget",  x + 1.44, y + h - 0.32, 1.2, 0.22, size=8, color=C_MUTED)

    def waterfall_chart(slide, data, x, y, w, h):
        """Positive = green bar, negative = red bar. Good for variance analysis."""
        if not data: return
        diffs = [float(d.get("Vs budget diff", d.get("variance", 0)) or 0) for d in data]
        max_v = max([abs(d) for d in diffs], default=1)
        if max_v == 0: max_v = 1

        n      = len(data)
        bar_w  = w / n
        area_h = h - 0.45
        mid_y  = y + area_h / 2

        add_rect(slide, x, mid_y, w, 0.01, C_FAINT)  # zero line

        for i, (d, diff) in enumerate(zip(data, diffs)):
            bx    = x + i * bar_w + bar_w * 0.15
            bw    = bar_w * 0.70
            color = C_GREEN if diff >= 0 else C_RED
            norm  = abs(diff) / max_v * (area_h / 2) * 0.9
            if diff >= 0:
                add_rect(slide, bx, mid_y - norm, bw, norm + 0.01, color)
            else:
                add_rect(slide, bx, mid_y, bw, norm, color)
            lbl = str(d.get("Label", d.get("Konto", "")))[:9]
            add_text(slide, lbl, bx - bar_w * 0.05, y + area_h + 0.04, bar_w, 0.24,
                     size=7, color=C_FAINT, align=PP_ALIGN.CENTER)

    def line_trend(slide, series, x, y, w, h, color=C_ACCENT, dashed_from=0):
        """Draw a simple connected dot-line using thin rectangles for trend."""
        if len(series) < 2: return
        max_v = max([abs(float(s.get("actual", s.get("forecast", 0)) or 0)) for s in series], default=1)
        if max_v == 0: max_v = 1
        area_h = h - 0.3
        pts = []
        for i, s in enumerate(series):
            v   = float(s.get("actual", s.get("forecast", 0)) or 0)
            px  = x + (i / (len(series) - 1)) * w
            py  = y + area_h - (abs(v) / max_v) * area_h * 0.88
            pts.append((px, py))
            # dot
            add_rect(slide, px - 0.05, py - 0.04, 0.10, 0.08, color)

        for i in range(len(pts) - 1):
            x1, y1 = pts[i]
            x2, y2 = pts[i + 1]
            # thin horizontal + vertical segments (L-shape approximation)
            mid = (x1 + x2) / 2
            add_rect(slide, x1, min(y1, y2) - 0.01, x2 - x1, abs(y2 - y1) + 0.02, color)

    def scenario_bars(slide, scenarios, x, y, w):
        labels = [("Optimistisk", C_GREEN), ("Bas", C_ACCENT), ("Pessimistisk", C_RED)]
        keys   = ["optimistic", "base", "pessimistic"]
        max_v  = max([scenarios.get(k, 1) for k in keys], default=1)
        if max_v == 0: max_v = 1
        for i, (k, (lbl, color)) in enumerate(zip(keys, labels)):
            val  = scenarios.get(k, 0)
            bar_w_px = (val / max_v) * w * 0.75
            row_y = y + i * 0.6
            add_rect(slide, x, row_y + 0.08, bar_w_px, 0.30, color)
            add_text(slide, lbl, x, row_y, 1.4, 0.26, size=9, color=C_MUTED)
            add_text(slide, fmt_sek(val), x + bar_w_px + 0.1, row_y + 0.06,
                     2.0, 0.26, size=10, bold=True, color=color)

    def variance_table(slide, rows, x, y, w, max_rows=8):
        """Compact variance table."""
        cols    = ["Konto", "Utfall", "Budget", "Avvikelse", "Avv %"]
        col_w   = [w * 0.32, w * 0.16, w * 0.16, w * 0.18, w * 0.18]
        row_h   = 0.265
        # Header
        cx = x
        for header, cw in zip(cols, col_w):
            add_text(slide, header, cx, y, cw - 0.04, 0.24, size=8,
                     color=C_FAINT, bold=True)
            cx += cw
        y += 0.26
        for row in rows[:max_rows]:
            add_rect(slide, x, y, w, row_h - 0.02, C_CARD)
            diff = float(row.get("Vs budget diff", 0))
            pct  = float(row.get("Vs budget %",   0))
            sign = C_GREEN if diff >= 0 else C_RED
            vals = [
                str(row.get("Label", row.get("Konto", "?")))[:28],
                fmt_sek(float(row.get("Utfall", row.get("actual", 0)))),
                fmt_sek(float(row.get("Budget", row.get("budget", 0)))),
                fmt_sek(diff),
                f"{pct*100:+.1f}%",
            ]
            colors = [C_MUTED, C_MUTED, C_MUTED, sign, sign]
            cx = x
            for val, cw, col in zip(vals, col_w, colors):
                add_text(slide, val, cx + 0.05, y + 0.04, cw - 0.09, 0.22,
                         size=8, color=col)
                cx += cw
            y += row_h
            if y > 6.8: break

    def insight_card(slide, x, y, w, h, title, items, accent=C_ACCENT):
        add_rect(slide, x, y, w, h, C_CARD)
        add_rect(slide, x, y, 0.04, h, accent)
        add_text(slide, title, x + 0.14, y + 0.10, w - 0.22, 0.26,
                 size=10, bold=True, color=C_ACCENT2)
        text_y = y + 0.38
        for item in items[:4]:
            if text_y + 0.22 > y + h - 0.08: break
            add_text(slide, f"• {item}", x + 0.14, text_y, w - 0.22, 0.24,
                     size=9, color=C_MUTED)
            text_y += 0.25

    # ════════════════════════════════════════════════════════════════
    # SLIDE BUILDERS  (one function per report type)
    # ════════════════════════════════════════════════════════════════

    def title_slide(subtitle=""):
        s = prs.slides.add_slide(blank); bg(s)
        add_rect(s, 0, 0, 0.07, 7.5, C_ACCENT)
        # Left accent panel
        add_rect(s, 0, 0, 4.2, 7.5, C_CARD)
        add_text(s, "NORDSHEET", 0.28, 0.38, 3.6, 0.35, size=10,
                 color=C_FAINT, bold=True)
        add_text(s, report_title, 0.28, 1.1, 3.65, 1.4, size=32, bold=True)
        add_text(s, subtitle or f"Period: {current_period}", 0.28, 2.7, 3.65, 0.4,
                 size=13, color=C_MUTED)
        add_text(s, f"Föregående period: {previous_period}", 0.28, 3.15, 3.65, 0.32,
                 size=10, color=C_FAINT)
        if context_text:
            add_text(s, f"Fokus: {context_text}", 0.28, 3.55, 3.65, 0.35,
                     size=10, color=C_ACCENT2, italic=True)
        # Right: hero KPIs
        kpi_box(s, 4.5, 1.3, 4.0, 1.4, "TOTALT UTFALL", fmt_sek(total_actual),
                val_color=C_WHITE)
        vc = C_GREEN if variance >= 0 else C_RED
        kpi_box(s, 4.5, 2.85, 4.0, 1.4, "AVVIKELSE VS BUDGET",
                f"{fmt_sek(variance)} ({var_pct*100:+.1f}%)", val_color=vc)
        mc = C_GREEN if mom_pct >= 0 else C_RED
        kpi_box(s, 4.5, 4.4, 4.0, 1.4, "MOM-FÖRÄNDRING",
                f"{mom_pct*100:+.1f}%", val_color=mc)
        # Narrative strip
        add_rect(s, 4.5, 5.95, 8.55, 1.3, C_CARD2)
        add_text(s, "AI-sammanfattning", 4.65, 5.98, 8.2, 0.28, size=9,
                 color=C_ACCENT2, bold=True)
        add_text(s, ai_exec_summ[:240], 4.65, 6.28, 8.2, 0.9, size=10, color=C_MUTED)
        return s

    def kpi_overview_slide():
        s = prs.slides.add_slide(blank); bg(s)
        slide_header(s, "Nyckeltal", f"Period {current_period}", report_title.upper())
        divider(s, 1.45)
        kpis = [
            ("UTFALL",          fmt_sek(total_actual),              "", C_WHITE),
            ("BUDGET",          fmt_sek(total_budget),              "", C_WHITE),
            ("AVVIKELSE",       fmt_sek(variance),                  f"{var_pct*100:+.1f}%",
             C_GREEN if variance >= 0 else C_RED),
            ("MOM",             f"{mom_pct*100:+.1f}%",             f"vs {previous_period}",
             C_GREEN if mom_pct >= 0 else C_RED),
        ]
        x0 = 0.28; spacing = 3.1
        for i, (lbl, val, sub, col) in enumerate(kpis):
            kpi_box(s, x0 + i * spacing, 1.65, 2.85, 1.5, lbl, val, sub, col)
        # Bar chart — period series
        add_text(s, "Utfall per period", 0.28, 3.38, 8, 0.3, size=11, bold=True, color=C_WHITE)
        bar_data = [{"label": p.get("period", ""), "actual": p.get("actual", 0),
                     "budget": p.get("budget", 0)} for p in period_series[-8:]]
        if bar_data:
            bar_chart(s, bar_data, 0.28, 3.72, 8.5, 3.2)
        # Right summary
        add_text(s, "Avvikelsetopp", 9.1, 3.38, 4.0, 0.3, size=11, bold=True, color=C_WHITE)
        for i, row in enumerate(top_budget[:4]):
            ry    = 3.72 + i * 0.72
            diff  = float(row.get("Vs budget diff", 0))
            color = C_GREEN if diff >= 0 else C_RED
            add_rect(s, 9.1, ry, 4.0, 0.62, C_CARD)
            add_text(s, str(row.get("Label", row.get("Konto", "?")))[:22],
                     9.24, ry + 0.04, 2.6, 0.28, size=9, color=C_MUTED)
            add_text(s, fmt_sek(diff), 9.24, ry + 0.32, 2.6, 0.24,
                     size=11, bold=True, color=color)
        return s

    def variance_detail_slide():
        s = prs.slides.add_slide(blank); bg(s)
        slide_header(s, "Avvikelseanalys", "Topp avvikelser mot budget", report_title.upper())
        divider(s, 1.45)
        # Waterfall chart
        add_text(s, "Budget vs utfall — avvikelse per konto",
                 0.28, 1.55, 8.5, 0.28, size=10, bold=True, color=C_WHITE)
        wf_data = (top_budget + top_mom)[:10]
        waterfall_chart(s, wf_data, 0.28, 1.88, 8.5, 3.2)
        # Table
        add_text(s, "Detaljerad avvikelsetabell",
                 0.28, 5.22, 12.7, 0.28, size=10, bold=True, color=C_WHITE)
        max_r = 6 if detail_level < 50 else 8
        variance_table(s, top_budget[:max_r], 0.28, 5.55, 12.7)
        return s

    def trend_slide():
        s = prs.slides.add_slide(blank); bg(s)
        slide_header(s, "Historisk trend", "Utfall per period", report_title.upper())
        divider(s, 1.45)
        if period_series:
            add_text(s, "Utfall & budget per period",
                     0.28, 1.58, 12.7, 0.3, size=11, bold=True, color=C_WHITE)
            bar_chart(s, period_series[-12:], 0.28, 1.92, 12.7, 5.2)
        else:
            add_text(s, "Ingen perioddata tillgänglig.", 0.28, 3.5, 12.7, 0.5,
                     size=13, color=C_FAINT, align=PP_ALIGN.CENTER)
        return s

    def forecast_slide():
        s = prs.slides.add_slide(blank); bg(s)
        slide_header(s, "Prognos", "Rullande 12-månadersprognos", report_title.upper())
        divider(s, 1.45)

        # Three scenario KPI boxes
        scenario_data = [
            ("OPTIMISTISK",  fc["scenarios"]["optimistic"], C_GREEN),
            ("BASSCENARIO",  fc["scenarios"]["base"],       C_ACCENT),
            ("PESSIMISTISK", fc["scenarios"]["pessimistic"],C_RED),
        ]
        for i, (lbl, val, col) in enumerate(scenario_data):
            kpi_box(s, 0.28 + i * 4.1, 1.65, 3.85, 1.35, lbl, fmt_sek(val),
                    val_color=col)

        # Scenario bar chart
        add_text(s, "Scenariojämförelse",
                 0.28, 3.18, 6.5, 0.3, size=11, bold=True, color=C_WHITE)
        scenario_bars(s, fc["scenarios"], 0.28, 3.52, 6.5)

        # Forecast mini bar chart
        add_text(s, "Prognos per period (F+1 → F+6)",
                 7.2, 3.18, 5.85, 0.3, size=11, bold=True, color=C_WHITE)
        fc_bars = [{"label": d["label"], "actual": d["forecast"]}
                   for d in fc["forecast"][:6]]
        if fc_bars:
            bar_chart(s, fc_bars, 7.2, 3.52, 5.85, 2.75, show_budget=False)

        # Additional KPIs
        kpi_box(s, 0.28,  5.6, 3.0, 1.25, "INKOMSTPROGNOS",
                fmt_sek(fc["income_forecast"]))
        kpi_box(s, 3.5,   5.6, 3.0, 1.25, "KOSTNADSPROGNOS",
                fmt_sek(fc["cost_forecast"]))
        kpi_box(s, 6.72,  5.6, 3.0, 1.25, "LÖNSAMHETSPROGNOS",
                f"{fc['margin_pct']:.1f}%",
                val_color=C_GREEN if fc["margin_pct"] > 0 else C_RED)
        kpi_box(s, 9.94,  5.6, 3.1, 1.25, "TILLVÄXTTAKT",
                f"{fc['growth_rate']:+.1f}%",
                val_color=C_GREEN if fc["growth_rate"] > 0 else C_RED)
        return s

    def insights_slide():
        s = prs.slides.add_slide(blank); bg(s)
        slide_header(s, "Insikter & Rekommendationer",
                     "AI-genererade observationer", report_title.upper())
        divider(s, 1.45)
        cols, card_w = 2, 6.1
        card_h = 2.4
        for i, insight in enumerate(ai_insights[:4]):
            col_i = i % cols; row_i = i // cols
            cx = 0.28 + col_i * (card_w + 0.44)
            cy = 1.58 + row_i * (card_h + 0.22)
            insight_card(s, cx, cy, card_w, card_h,
                         insight["title"], insight.get("items", []))
        return s

    def action_items_slide(items):
        s = prs.slides.add_slide(blank); bg(s)
        slide_header(s, "Åtgärdspunkter & Nästa steg",
                     f"Prioriteringar efter {current_period}", report_title.upper())
        divider(s, 1.45)
        for i, item in enumerate(items[:6]):
            ry = 1.65 + i * 0.82
            add_rect(s, 0.28, ry, 12.7, 0.72, C_CARD)
            add_rect(s, 0.28, ry, 0.36, 0.72, C_ACCENT)
            add_text(s, str(i + 1), 0.35, ry + 0.18, 0.22, 0.36,
                     size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
            add_text(s, str(item), 0.78, ry + 0.18, 11.9, 0.38,
                     size=11, color=C_MUTED)
        return s

    def end_slide():
        s = prs.slides.add_slide(blank); bg(s)
        add_rect(s, 0, 0, 0.07, 7.5, C_ACCENT)
        add_text(s, "NORDSHEET", 0.25, 3.0, 12.8, 0.45,
                 size=11, color=C_FAINT, bold=True, align=PP_ALIGN.CENTER)
        add_text(s,
                 "Finance intelligence for the modern controller.",
                 0.25, 3.55, 12.8, 0.55, size=18, color=C_MUTED,
                 align=PP_ALIGN.CENTER)
        add_text(s, f"Rapport: {report_title} · Period: {current_period}",
                 0.25, 4.2, 12.8, 0.35, size=10, color=C_FAINT,
                 align=PP_ALIGN.CENTER)
        return s

    # ════════════════════════════════════════════════════════════════
    # COMPOSE DECKS PER REPORT TYPE
    # ════════════════════════════════════════════════════════════════

    action_list = []
    for g in ai_insights:
        if "action" in g.get("title", "").lower() or "åtgärd" in g.get("title", "").lower():
            action_list = g.get("items", [])
            break
    if not action_list:
        action_list = [
            f"Granska avvikelse för {top_budget[0].get('Label', '?')}" if top_budget else "Granska avvikelser",
            "Uppdatera prognos baserat på utfall",
            "Kommunicera åtgärdsplan till ansvariga",
        ]

    if report_type in ("monthly", "income_statement", "budget_vs_actual"):
        title_slide(f"Månadsavstämning — {current_period}")
        kpi_overview_slide()
        variance_detail_slide()
        if detail_level >= 40 and period_series:
            trend_slide()
        forecast_slide()
        insights_slide()
        action_items_slide(action_list)

    elif report_type == "quarterly":
        title_slide(f"Kvartal — {current_period}")
        kpi_overview_slide()
        trend_slide()
        variance_detail_slide()
        forecast_slide()
        insights_slide()
        action_items_slide(action_list)

    elif report_type in ("annual", "balance_sheet", "cash_flow"):
        title_slide(f"Helårsanalys — {current_period}")
        kpi_overview_slide()
        trend_slide()
        variance_detail_slide()
        forecast_slide()
        insights_slide()
        # Extra: full variance table slide for annual
        s = prs.slides.add_slide(blank); bg(s)
        slide_header(s, "Komplett avvikelseregister", "Alla konton", report_title.upper())
        divider(s, 1.45)
        all_rows = top_budget + [r for r in top_mom if r not in top_budget]
        variance_table(s, all_rows, 0.28, 1.58, 12.7, max_rows=14)
        action_items_slide(action_list)

    elif report_type == "forecast":
        title_slide("Helårsprognos & Scenarier")
        # Forecast is the hero
        forecast_slide()
        kpi_overview_slide()
        if period_series:
            trend_slide()
        insights_slide()
        action_items_slide(action_list)

    elif report_type == "kpi_dashboard":
        title_slide("KPI-rapport")
        kpi_overview_slide()
        # KPI cards 2×2 layout slide
        s = prs.slides.add_slide(blank); bg(s)
        slide_header(s, "KPI-dashboard", f"Status {current_period}", report_title.upper())
        divider(s, 1.45)
        kpi_rows = [
            ("UTFALL",          fmt_sek(total_actual),     f"Budget: {fmt_sek(total_budget)}", C_WHITE),
            ("AVVIKELSE",       fmt_sek(variance),         f"{var_pct*100:+.1f}% vs budget",
             C_GREEN if variance >= 0 else C_RED),
            ("MOM-TREND",       f"{mom_pct*100:+.1f}%",   f"vs {previous_period}",
             C_GREEN if mom_pct >= 0 else C_RED),
            ("PROGNOS (BAS)",   fmt_sek(fc["scenarios"]["base"]), "Nästa period", C_ACCENT2),
            ("INKOMST FC",      fmt_sek(fc["income_forecast"]), "", C_GREEN),
            ("KOSTNAD FC",      fmt_sek(fc["cost_forecast"]),   "", C_RED),
        ]
        positions = [(0.28,1.65),(4.42,1.65),(8.56,1.65),(0.28,3.25),(4.42,3.25),(8.56,3.25)]
        for (kx, ky), (lbl, val, sub, col) in zip(positions, kpi_rows):
            kpi_box(s, kx, ky, 3.85, 1.35, lbl, val, sub, col)
        # Trend bar
        if period_series:
            add_text(s, "Trendöversikt", 0.28, 4.72, 12.7, 0.3,
                     size=11, bold=True, color=C_WHITE)
            bar_chart(s, period_series[-6:], 0.28, 5.05, 12.7, 2.2)
        insights_slide()
        action_items_slide(action_list)

    else:
        # Fallback: generic deck
        title_slide()
        kpi_overview_slide()
        variance_detail_slide()
        forecast_slide()
        insights_slide()
        action_items_slide(action_list)

    end_slide()

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()



# ═══════════════════════════════════════════════════════════════════
# DOCX BUILDER
# ═══════════════════════════════════════════════════════════════════

def build_docx(pack: dict, spec: dict, ai_context: str = "") -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()

    # Page margins
    for section in doc.sections:
        section.top_margin    = Cm(2.0)
        section.bottom_margin = Cm(2.0)
        section.left_margin   = Cm(2.5)
        section.right_margin  = Cm(2.5)

    # Data
    current_period  = pack.get("current_period",  "—")
    previous_period = pack.get("previous_period", "—")
    total_actual    = float(pack.get("total_actual", 0))
    total_budget    = float(pack.get("total_budget", 0))
    variance        = total_actual - total_budget
    var_pct         = safe_pct(total_actual, total_budget) or 0
    kpi             = pack.get("kpi_summary", [{}])[0]
    mom_pct         = float(kpi.get("MoM %") or 0)
    narrative       = pack.get("narrative", "")
    top_budget      = pack.get("top_budget", [])
    top_mom         = pack.get("top_mom",    [])
    report_type     = spec.get("report_type",  "monthly")
    tone            = spec.get("tone",         "professional")
    context_text    = spec.get("context",      "") or ai_context

    REPORT_TITLES = {
        "income_statement": "Resultaträkning",
        "balance_sheet":    "Balansräkning",
        "cash_flow":        "Kassaflödesanalys",
        "budget_vs_actual": "Budget vs Utfall",
        "ai_summary":       "AI-genererad sammanfattning",
        "monthly":          "Månadsrapport",
        "quarterly":        "Kvartalsrapport",
        "annual":           "Årsredovisning",
    }
    report_title = spec.get("title") or REPORT_TITLES.get(report_type, "Finansiell rapport")

    # AI narrative
    openai_key = os.getenv("OPENAI_API_KEY")
    full_narrative = narrative
    executive_summary = ""

    if openai_key:
        try:
            from openai import OpenAI
            client = OpenAI(api_key=openai_key)
            prompt = (
                f"Du är en senior finansanalytiker. Skriv en {report_title} på svenska med ton: {tone}. "
                f"Inkludera: executive summary (2 stycken), analys av avvikelser, "
                f"rekommendationer och slutsats. "
                f"Data: Utfall {fmt_sek(total_actual)}, Budget {fmt_sek(total_budget)}, "
                f"Avvikelse {fmt_sek(variance)} ({var_pct*100:+.1f}%), MoM {mom_pct*100:+.1f}%. "
                f"Topp-avvikelser: {', '.join([r.get('Label','') + ' ' + fmt_sek(r.get('Vs budget diff',0)) for r in top_budget[:5]])}. "
                + (f"Extra instruktion: {context_text}" if context_text else "")
                + " Returnera JSON: {\"executive_summary\": \"...\", \"analysis\": \"...\", \"recommendations\": \"...\", \"conclusion\": \"...\"}"
            )
            resp = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "user", "content": prompt}],
                max_tokens=1000,
            )
            raw = re.sub(r"```json|```", "", resp.choices[0].message.content or "{}").strip()
            ai_text = json.loads(raw)
            executive_summary = ai_text.get("executive_summary", "")
            full_narrative     = ai_text.get("analysis", narrative)
        except Exception:
            executive_summary = narrative

    def h1(text):
        p = doc.add_heading(text, level=1)
        p.style.font.color.rgb = RGBColor(0x6C, 0x63, 0xFF)

    def h2(text):
        doc.add_heading(text, level=2)

    def para(text, italic=False):
        p = doc.add_paragraph(text)
        if italic:
            for run in p.runs:
                run.italic = True
        return p

    def kv(label, value, color=None):
        p = doc.add_paragraph()
        r1 = p.add_run(f"{label}: ")
        r1.bold = True
        r2 = p.add_run(value)
        if color:
            r2.font.color.rgb = color

    # ── Cover ──
    doc.add_heading("NORDSHEET", 0)
    doc.add_heading(report_title, 0)
    para(f"Period: {current_period}")
    if previous_period:
        para(f"Jämförs med: {previous_period}")
    if context_text:
        para(f"Fokus: {context_text}", italic=True)
    doc.add_page_break()

    # ── Executive Summary ──
    h1("Executive Summary")
    para(executive_summary or narrative)
    doc.add_paragraph()

    # ── KPI ──
    h1("Nyckeltal")
    kv("Totalt utfall",   fmt_sek(total_actual))
    kv("Budget",          fmt_sek(total_budget))
    kv("Avvikelse",       f"{fmt_sek(variance)} ({var_pct*100:+.1f}%)",
       RGBColor(0x22,0xC5,0x5E) if variance >= 0 else RGBColor(0xEF,0x44,0x44))
    kv("MoM-förändring",  f"{mom_pct*100:+.1f}% vs {previous_period or '—'}",
       RGBColor(0x22,0xC5,0x5E) if mom_pct >= 0 else RGBColor(0xEF,0x44,0x44))

    doc.add_paragraph()

    # ── Analysis ──
    h1("Analys")
    para(full_narrative)

    # ── Variance Table ──
    h1("Avvikelseanalys")
    h2("Negativa avvikelser (överskridningar)")

    if top_budget:
        t = doc.add_table(rows=1, cols=5)
        t.style = "Table Grid"
        hdr = t.rows[0].cells
        for i, col in enumerate(["Konto", "Utfall", "Budget", "Avvikelse", "Avv %"]):
            hdr[i].text = col
            hdr[i].paragraphs[0].runs[0].bold = True

        for row in top_budget[:10]:
            cells = t.add_row().cells
            cells[0].text = str(row.get("Label") or row.get("Konto") or "")
            cells[1].text = fmt_sek(float(row.get("Utfall", 0)))
            cells[2].text = fmt_sek(float(row.get("Budget", 0)))
            cells[3].text = fmt_sek(float(row.get("Vs budget diff", 0)))
            cells[4].text = f"{float(row.get('Vs budget %', 0))*100:+.1f}%"

    doc.add_paragraph()
    h2("Positiva avvikelser (besparingar)")

    if top_mom:
        t2 = doc.add_table(rows=1, cols=5)
        t2.style = "Table Grid"
        hdr2 = t2.rows[0].cells
        for i, col in enumerate(["Konto", "Utfall", "Budget", "Avvikelse", "Avv %"]):
            hdr2[i].text = col
            hdr2[i].paragraphs[0].runs[0].bold = True
        for row in top_mom[:10]:
            cells = t2.add_row().cells
            cells[0].text = str(row.get("Label") or row.get("Konto") or "")
            cells[1].text = fmt_sek(float(row.get("Utfall", 0)))
            cells[2].text = fmt_sek(float(row.get("Budget", 0)))
            cells[3].text = fmt_sek(float(row.get("Vs budget diff", 0)))
            cells[4].text = f"{float(row.get('Vs budget %', 0))*100:+.1f}%"

    # ── Forecast ──
    h1("Prognos")
    fc = build_forecast_data(pack)
    kv("Inkomstprognos",   fmt_sek(fc["income_forecast"]))
    kv("Kostnadsprognos",  fmt_sek(fc["cost_forecast"]))
    kv("Lönsamhetsprognos", f"{fc['margin_pct']:.1f}%")
    doc.add_paragraph()
    kv("Optimistiskt scenario",  fmt_sek(fc["scenarios"]["optimistic"]))
    kv("Basscenario",            fmt_sek(fc["scenarios"]["base"]))
    kv("Pessimistiskt scenario", fmt_sek(fc["scenarios"]["pessimistic"]))

    # ── Recommendations ──
    if openai_key:
        try:
            from openai import OpenAI
            client = OpenAI(api_key=openai_key)
            rp = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "user", "content":
                    f"Ge 3-5 konkreta rekommendationer på svenska för en {report_title}. "
                    f"Data: avvikelse {fmt_sek(variance)}, MoM {mom_pct*100:+.1f}%. "
                    f"Returnera JSON: [\"...\", \"...\"]"}],
                max_tokens=400,
            )
            raw_r = re.sub(r"```json|```", "", rp.choices[0].message.content or "[]").strip()
            recs  = json.loads(raw_r)
            h1("Rekommendationer")
            for rec in recs:
                doc.add_paragraph(f"• {rec}")
        except Exception:
            pass

    # ── Conclusion ──
    h1("Slutsats")
    para(
        f"Sammanfattningsvis visar {current_period} ett utfall på {fmt_sek(total_actual)} "
        f"mot budget {fmt_sek(total_budget)}, en avvikelse på {fmt_sek(variance)} ({var_pct*100:+.1f}%). "
        f"MoM-förändringen är {mom_pct*100:+.1f}%."
    )

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


# ═══════════════════════════════════════════════════════════════════
# ROUTES
# ═══════════════════════════════════════════════════════════════════

@app.get("/api/health")
def health():
    return {"ok": True}


@app.post("/api/analyze")
async def analyze(file: UploadFile = File(...)):
    try:
        contents = await file.read()
        df = read_upload_bytes(contents, file.filename or "")
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Kunde inte läsa filen: {e}")
    cols = list(df.columns)
    return {
        "available_columns": cols,
        "column_suggestions": suggest_columns(cols),
        "row_count": len(df),
        "preview": df.head(20).fillna("").to_dict(orient="records"),
    }


@app.post("/api/analyze-with-mapping")
async def analyze_with_mapping(
    file: UploadFile = File(...),
    mapping_json: str = Form(...),
):
    try:
        contents = await file.read()
        df = read_upload_bytes(contents, file.filename or "")
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Kunde inte läsa filen: {e}")
    try:
        mapping = json.loads(mapping_json)
    except Exception:
        raise HTTPException(status_code=422, detail="Ogiltig mapping JSON")
    try:
        pack = compute_pack(df, mapping)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Analysen misslyckades: {e}")
    return pack


class AiMapRequest(BaseModel):
    preview: dict        # { "ColName": {"dtype": "string", "examples": ["v1","v2"]} }
    business_type: str = ""


@app.post("/api/ai-map-columns")
async def ai_map_columns(req: AiMapRequest):
    """
    Tar emot ett kolumnpreview och returnerar AI-förslag på fältmappning
    via GPT-4o-mini. OPENAI_API_KEY hämtas från Vercel environment variables.

    Anropas från frontend ConnectPage som komplement till den regelbaserade
    suggest_columns() — ger bättre träff på icke-standardiserade kolumnnamn.
    """
    openai_key = os.getenv("OPENAI_API_KEY")
    if not openai_key:
        raise HTTPException(status_code=503, detail="OPENAI_API_KEY saknas i miljövariablerna.")

    target_fields = [
        "period", "account", "actual", "budget",
        "entity", "cost_center", "project", "account_name",
    ]

    prompt = (
        "Du är en senior business controller och expert på att tolka Excel-filer.\n\n"
        "Identifiera vilka kolumner som bäst motsvarar dessa fält:\n"
        + "\n".join(f"- {f}" for f in target_fields)
        + f"\n\nBolagstyp: {req.business_type or 'Okänd'}\n\n"
        "Kolumner med exempelvärden:\n"
        + json.dumps(req.preview, ensure_ascii=False, indent=2)
        + "\n\nReturnera ENDAST giltig JSON utan förklaringar eller backticks. "
        f"Nycklar: {', '.join(target_fields)}. "
        "Värde = kolumnnamnet vid match, annars null."
    )

    try:
        from openai import OpenAI
        client = OpenAI(api_key=openai_key)
        resp = client.chat.completions.create(
            model="gpt-4o-mini",
            temperature=0,
            max_tokens=300,
            messages=[
                {"role": "system", "content": "Du är expert på finansiell datamappning. Returnera ALLTID bara JSON."},
                {"role": "user",   "content": prompt},
            ],
        )
        raw = re.sub(r"```json|```", "", resp.choices[0].message.content or "{}").strip()
        try:
            parsed = json.loads(raw)
        except Exception:
            m = re.search(r"\{.*\}", raw, re.DOTALL)
            parsed = json.loads(m.group(0)) if m else {}

        valid_cols = set(req.preview.keys())
        return {
            f: (parsed[f] if isinstance(parsed.get(f), str) and parsed[f] in valid_cols else None)
            for f in target_fields
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI mapping misslyckades: {e}")


class ChatRequest(BaseModel):
    question: str
    pack: Any


@app.post("/api/chat")
async def chat(req: ChatRequest):
    openai_key = os.getenv("OPENAI_API_KEY")
    pack = req.pack or {}

    # Detect if this is a variance analysis request (has structured 4-point format)
    is_variance_analysis = all(k in req.question for k in ["ORSAK", "TYP", "ÅTGÄRD", "PROGNOS"])

    if openai_key:
        try:
            from openai import OpenAI
            client = OpenAI(api_key=openai_key)

            if is_variance_analysis:
                # Deep variance analysis — structured, specific, actionable
                system = """Du är en senior finanscontroller med 15 års erfarenhet.
Du analyserar avvikelser med precision och ger konkreta, handlingsbara svar.

Regler:
- Svara ALLTID på svenska
- Svara strukturerat med exakt de rubriker användaren angett (1. ORSAK, 2. TYP, 3. ÅTGÄRD, 4. PROGNOS)
- Var specifik — undvik generella fraser som "kan bero på"
- Max 2 meningar per punkt
- För TYP: basera på om mönstret är synligt i historiken
- För PROGNOS: räkna ut helårssiffran om avvikelsen är återkommande
- Använd siffror från datan, inte vaga uppskattningar"""

                # Build context from pack
                period_series = pack.get("period_series", [])
                history_str = ""
                if period_series:
                    history_str = "\n".join([
                        f"  {p.get('period','?')}: utfall {p.get('actual',0):,.0f}, budget {p.get('budget',0):,.0f}"
                        for p in period_series[-6:]
                    ])

                context = f"""Finansiell kontext:
- Aktuell period: {pack.get('current_period', '?')}
- Föregående period: {pack.get('previous_period', '?')}
- Totalt utfall: {pack.get('total_actual', 0):,.0f}
- Total budget: {pack.get('total_budget', 0):,.0f}
{f'Periodhistorik (senaste 6):{chr(10)}{history_str}' if history_str else ''}"""

                resp = client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[
                        {"role": "system",  "content": system},
                        {"role": "user",    "content": f"{context}\n\n{req.question}"},
                    ],
                    max_tokens=600,
                    temperature=0.1,
                )
            else:
                # General chat about financials
                system = (
                    "Du är en senior finansanalytiker på NordSheet. "
                    "Svara kort, konkret och på svenska. Använd siffror från datan. "
                    f"Finansdata: period={pack.get('current_period','?')}, "
                    f"utfall={pack.get('total_actual',0):,.0f}, "
                    f"budget={pack.get('total_budget',0):,.0f}"
                )
                resp = client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[
                        {"role": "system", "content": system},
                        {"role": "user",   "content": req.question},
                    ],
                    max_tokens=400,
                    temperature=0.2,
                )

            return {"answer": resp.choices[0].message.content}
        except Exception as e:
            return {"answer": f"AI-fel: {e}"}

    # Fallback without OpenAI key
    q = req.question.lower()
    if any(w in q for w in ["utfall", "actual", "total"]):
        return {"answer": f"Totalt utfall: {fmt_sek(float(pack.get('total_actual', 0)))}"}
    if "budget" in q:
        return {"answer": f"Total budget: {fmt_sek(float(pack.get('total_budget', 0)))}"}
    if "period" in q:
        return {"answer": f"Aktuell period: {pack.get('current_period', 'N/A')}"}
    return {"answer": pack.get("narrative", "Ingen data tillgänglig.")}


class ExportRequest(BaseModel):
    fmt:           str
    pack:          Any
    report_items:  list = []
    spec:          dict = {}
    purpose:       str  = ""
    business_type: str  = ""
"""
Fakturaanalys-endpoint.
Klistra in detta block i main.py, direkt efter /api/chat-endpointen.

Kräver att 'pypdf' eller 'pdfplumber' finns i requirements.txt.
Lägg till: pdfplumber>=0.10.3
"""


@app.post("/api/analyze-invoice")
async def analyze_invoice(file: UploadFile = File(...)):
    """
    Tar emot en PDF-faktura, extraherar text och skickar till OpenAI
    för strukturerad analys. Returnerar leverantör, belopp, rader och
    AI-flaggade avvikelser.
    """
    if not file.filename or not file.filename.lower().endswith(".pdf"):
        raise HTTPException(status_code=422, detail="Endast PDF-filer stöds.")

    try:
        contents = await file.read()
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Kunde inte läsa filen: {e}")

    # ── Extrahera text från PDF ──
    pdf_text = ""
    try:
        import pdfplumber, io as _io
        with pdfplumber.open(_io.BytesIO(contents)) as pdf:
            pages = []
            for page in pdf.pages[:6]:  # max 6 sidor
                t = page.extract_text()
                if t:
                    pages.append(t)
            pdf_text = "\n\n".join(pages)
    except ImportError:
        # Fallback: pypdf
        try:
            from pypdf import PdfReader
            import io as _io
            reader = PdfReader(_io.BytesIO(contents))
            pdf_text = "\n\n".join(
                page.extract_text() or ""
                for page in reader.pages[:6]
            )
        except ImportError:
            raise HTTPException(
                status_code=500,
                detail="pdfplumber eller pypdf saknas. Lägg till i requirements.txt."
            )
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"Kunde inte läsa PDF: {e}")
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"PDF-extraktion misslyckades: {e}")

    if not pdf_text.strip():
        raise HTTPException(
            status_code=422,
            detail="Kunde inte extrahera text från PDF. Filen kan vara skannad eller skyddad."
        )

    # Begränsa textstorlek
    pdf_text = pdf_text[:6000]

    # ── AI-analys ──
    openai_key = os.getenv("OPENAI_API_KEY")

    fallback_result = {
        "supplier":       "Okänd",
        "invoice_number": "",
        "invoice_date":   "",
        "due_date":       "",
        "total_amount":   None,
        "vat_amount":     None,
        "net_amount":     None,
        "currency":       "SEK",
        "line_items":     [],
        "ai_summary":     "AI-analys ej tillgänglig — OPENAI_API_KEY saknas.",
        "anomalies":      [],
        "category":       "Okänd",
        "confidence":     0.5,
    }

    if not openai_key:
        return fallback_result

    try:
        from openai import OpenAI
        client = OpenAI(api_key=openai_key)

        prompt = f"""Du är en expert på fakturaanalys. Analysera denna faktura och returnera ENDAST giltig JSON.



Fakturatextinnehåll:
---
{pdf_text}
---

Returnera JSON med exakt dessa fält:
{{
  "supplier": "leverantörens namn",
  "invoice_number": "fakturanummer",
  "invoice_date": "YYYY-MM-DD eller tom sträng",
  "due_date": "YYYY-MM-DD eller tom sträng",
  "total_amount": 1234.56 (number eller null),
  "vat_amount": 123.45 (number eller null),
  "net_amount": 1111.11 (number eller null),
  "currency": "SEK",
  "category": "t.ex. Konsulttjänster / IT / Hyra / Transport / Övrigt",
  "line_items": [
    {{"description": "...", "amount": 100.0, "quantity": 1}}
  ],
  "anomalies": [
    "Beskriv avvikelser här t.ex. saknat org.nr, ovanligt högt belopp, förfallet datum"
  ],
  "ai_summary": "2-3 meningar på svenska om fakturan — vad den avser, om den ser korrekt ut och vad controllern bör tänka på",
  "confidence": 0.85
}}

Regler:
- Alla belopp som number (inte sträng)
- Datum alltid YYYY-MM-DD format
- anomalies: lista verkliga avvikelser du hittar — saknat org.nr, förfallodatum passerat, moms ser fel ut, dublett-risk etc.
- Om du inte kan hitta ett värde, använd null för tal och tom sträng för text
- Returnera BARA JSON, inga förklaringar"""

        resp = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "Du är expert på fakturaanalys. Returnera alltid giltig JSON."},
                {"role": "user",   "content": prompt},
            ],
            max_tokens=1000,
            temperature=0.1,
        )

        raw = resp.choices[0].message.content or "{}"
        raw = re.sub(r"```json|```", "", raw).strip()

        try:
            result = json.loads(raw)
        except Exception:
            m = re.search(r"\{.*\}", raw, re.DOTALL)
            result = json.loads(m.group(0)) if m else {}

        # Säkerställ att alla fält finns
        for key, default in fallback_result.items():
            if key not in result:
                result[key] = default

        # Säkerställ att numeriska fält verkligen är numbers eller null
        for num_field in ["total_amount", "vat_amount", "net_amount"]:
            val = result.get(num_field)
            if val is not None:
                try:
                    result[num_field] = float(val)
                except (TypeError, ValueError):
                    result[num_field] = None

        return result

    except Exception as e:
        fallback_result["ai_summary"] = f"AI-analys misslyckades: {str(e)}"
        return fallback_result
@app.post("/api/invoice-inbound")
async def invoice_inbound(request: Request):
    """
    Postmark skickar hit när ett mail kommer in till er inbound-adress.
    Extraherar PDF-bilagor och kör fakturaanalys automatiskt.
    """
    try:
        data = await request.json()
    except Exception:
        return {"ok": False, "error": "invalid json"}

    attachments = data.get("Attachments", [])
    from_email  = data.get("From", "")
    subject     = data.get("Subject", "")

    results = []

    for attachment in attachments:
        filename    = attachment.get("Name", "")
        content_b64 = attachment.get("Content", "")
        content_type = attachment.get("ContentType", "")

        # Bara PDF-bilagor
        if not (filename.lower().endswith(".pdf") or "pdf" in content_type.lower()):
            continue

        try:
            import base64, io as _io
            pdf_bytes = base64.b64decode(content_b64)

            # Extrahera text
            pdf_text = ""
            try:
                import pdfplumber
                with pdfplumber.open(_io.BytesIO(pdf_bytes)) as pdf:
                    pdf_text = "\n\n".join(
                        p.extract_text() or ""
                        for p in pdf.pages[:6]
                    )
            except Exception:
                pass

            if not pdf_text.strip():
                results.append({"file": filename, "error": "Ingen text i PDF"})
                continue

            # Kör AI-analys (återanvänd samma logik)
            openai_key = os.getenv("OPENAI_API_KEY")
            if openai_key:
                from openai import OpenAI
                client = OpenAI(api_key=openai_key)
                prompt = f"""Analysera denna faktura och returnera JSON:
{{
  "supplier": "...",
  "invoice_number": "...",
  "invoice_date": "YYYY-MM-DD",
  "due_date": "YYYY-MM-DD",
  "total_amount": 0.0,
  "vat_amount": 0.0,
  "net_amount": 0.0,
  "currency": "SEK",
  "category": "...",
  "line_items": [],
  "anomalies": [],
  "ai_summary": "...",
  "confidence": 0.9
}}

Faktura:
{pdf_text[:4000]}"""

                resp = client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[{"role": "user", "content": prompt}],
                    max_tokens=800,
                    temperature=0.1,
                )
                raw = re.sub(r"```json|```", "",
                    resp.choices[0].message.content or "{}").strip()
                analysis = json.loads(raw)
                results.append({
                    "file":     filename,
                    "from":     from_email,
                    "subject":  subject,
                    "analysis": analysis,
                })

        except Exception as e:
            results.append({"file": filename, "error": str(e)})

    logger.info(f"Inbound invoice from {from_email}: {len(results)} PDFs processed")
    return {"ok": True, "processed": len(results), "results": results}

@app.post("/api/export")
async def export_file(req: ExportRequest):
    pack        = req.pack or {}
    spec        = req.spec or {}
    ai_context  = spec.get("context", "")

    if req.fmt == "pptx":
        try:
            data = build_pptx(pack, spec, ai_context)
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"PPTX-generering misslyckades: {e}")
        return StreamingResponse(
            io.BytesIO(data),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": 'attachment; filename="nordsheet_rapport.pptx"'},
        )

    elif req.fmt == "docx":
        try:
            data = build_docx(pack, spec, ai_context)
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"DOCX-generering misslyckades: {e}")
        return StreamingResponse(
            io.BytesIO(data),
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": 'attachment; filename="nordsheet_rapport.docx"'},
        )

    raise HTTPException(status_code=400, detail=f"Okänt format: {req.fmt}")


@app.post("/api/forecast")
async def forecast_api(req: ChatRequest):
    """Return forecast data computed from pack."""
    try:
        fc = build_forecast_data(req.pack or {})
        return fc
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/forecast")
async def forecast_get():
    return {"message": "POST to /api/forecast with pack data."}
